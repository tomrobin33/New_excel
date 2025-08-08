"""
文档提取模块 - 从PPT、Word、PDF文件中提取表格数据
支持从各种文档格式中识别和提取表格，并转换为Excel格式
"""

import logging
import os
import tempfile
import uuid
from typing import List, Dict, Any, Optional, Union
import requests
import pandas as pd
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows

# 导入文档处理库
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available, PPT processing disabled")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available, Word processing disabled")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logging.warning("pdfplumber not available, PDF processing disabled")

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    logging.warning("PyPDF2 not available, PDF processing disabled")

logger = logging.getLogger(__name__)

class DocumentExtractor:
    """文档提取器类，用于从各种文档格式中提取表格数据"""
    
    def __init__(self):
        self.supported_formats = {
            'pptx': PPTX_AVAILABLE,
            'ppt': PPTX_AVAILABLE,
            'docx': DOCX_AVAILABLE,
            'doc': DOCX_AVAILABLE,
            'pdf': PDFPLUMBER_AVAILABLE or PYPDF2_AVAILABLE
        }
    
    def download_file(self, url: str) -> str:
        """从URL下载文件到临时目录"""
        try:
            # 确保URL正确编码
            from urllib.parse import urlparse, parse_qs, urlencode, urlunparse
            
            # 解析URL并重新编码
            parsed = urlparse(url)
            query_params = parse_qs(parsed.query)
            
            # 重新构建URL，确保参数正确编码
            encoded_query = urlencode(query_params, doseq=True)
            clean_url = urlunparse((
                parsed.scheme,
                parsed.netloc,
                parsed.path,
                parsed.params,
                encoded_query,
                parsed.fragment
            ))
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': '*/*',
                'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1'
            }
            
            logger.info(f"尝试下载文件: {clean_url}")
            response = requests.get(clean_url, stream=True, headers=headers, timeout=30, allow_redirects=True)
            
            # 详细的错误信息
            if response.status_code != 200:
                logger.error(f"HTTP错误: {response.status_code} - {response.reason}")
                logger.error(f"响应头: {dict(response.headers)}")
                if response.text:
                    logger.error(f"响应内容: {response.text[:500]}")
                raise Exception(f"HTTP {response.status_code}: {response.reason}")
            
            response.raise_for_status()
            
            # 获取文件扩展名
            content_disposition = response.headers.get('content-disposition', '')
            if 'filename=' in content_disposition:
                filename = content_disposition.split('filename=')[1].strip('"')
            else:
                # 从URL中提取文件名，去除查询参数
                filename = url.split('/')[-1].split('?')[0]
            
            # 确保文件名不为空且有效
            if not filename or len(filename) > 100:
                # 如果文件名过长或为空，使用默认名称
                filename = "document"
            
            # 添加文件扩展名（如果缺失）
            if '.' not in filename:
                content_type = response.headers.get('content-type', '').lower()
                if 'powerpoint' in content_type or 'presentation' in content_type:
                    filename += '.pptx'
                elif 'word' in content_type or 'document' in content_type:
                    filename += '.docx'
                elif 'pdf' in content_type:
                    filename += '.pdf'
                else:
                    filename += '.tmp'
            
            # 创建临时文件，使用短文件名
            temp_file = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4().hex[:8]}_{filename}")
            
            with open(temp_file, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
            
            return temp_file
            
        except Exception as e:
            logger.error(f"下载文件失败: {e}")
            raise Exception(f"无法下载文件: {str(e)}")
    
    def extract_tables_from_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """从PPTX文件中提取表格"""
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx库未安装，无法处理PPT文件")
        
        tables = []
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape_num, shape in enumerate(slide.shapes, 1):
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        
                        if table_data:
                            tables.append({
                                'slide': slide_num,
                                'shape': shape_num,
                                'data': table_data,
                                'rows': len(table_data),
                                'columns': len(table_data[0]) if table_data else 0
                            })
            
            return tables
            
        except Exception as e:
            logger.error(f"处理PPTX文件失败: {e}")
            raise Exception(f"处理PPTX文件失败: {str(e)}")
    
    def extract_tables_from_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """从DOCX文件中提取表格"""
        if not DOCX_AVAILABLE:
            raise Exception("python-docx库未安装，无法处理Word文件")
        
        tables = []
        try:
            doc = Document(file_path)
            
            for table_num, table in enumerate(doc.tables, 1):
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                if table_data:
                    tables.append({
                        'table': table_num,
                        'data': table_data,
                        'rows': len(table_data),
                        'columns': len(table_data[0]) if table_data else 0
                    })
            
            return tables
            
        except Exception as e:
            logger.error(f"处理DOCX文件失败: {e}")
            raise Exception(f"处理DOCX文件失败: {str(e)}")
    
    def extract_tables_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """从PDF文件中提取表格"""
        if not (PDFPLUMBER_AVAILABLE or PYPDF2_AVAILABLE):
            raise Exception("pdfplumber或PyPDF2库未安装，无法处理PDF文件")
        
        tables = []
        try:
            # 优先使用pdfplumber，效果更好
            if PDFPLUMBER_AVAILABLE:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        page_tables = page.extract_tables()
                        for table_num, table in enumerate(page_tables, 1):
                            if table:
                                # 清理表格数据
                                cleaned_table = []
                                for row in table:
                                    cleaned_row = []
                                    for cell in row:
                                        if cell is None:
                                            cleaned_row.append("")
                                        else:
                                            cleaned_row.append(str(cell).strip())
                                    cleaned_table.append(cleaned_row)
                                
                                if cleaned_table:
                                    tables.append({
                                        'page': page_num,
                                        'table': table_num,
                                        'data': cleaned_table,
                                        'rows': len(cleaned_table),
                                        'columns': len(cleaned_table[0]) if cleaned_table else 0
                                    })
            
            # 如果没有找到表格，尝试使用PyPDF2（作为备用方案）
            elif PYPDF2_AVAILABLE and not tables:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    # PyPDF2主要用于文本提取，表格提取能力有限
                    # 这里可以添加基本的文本表格识别逻辑
                    pass
            
            return tables
            
        except Exception as e:
            logger.error(f"处理PDF文件失败: {e}")
            raise Exception(f"处理PDF文件失败: {str(e)}")
    
    def extract_tables_from_document(self, file_path: str, file_type: str = None) -> List[Dict[str, Any]]:
        """从文档中提取表格，自动识别文件类型"""
        if not file_type:
            file_type = os.path.splitext(file_path)[1].lower().lstrip('.')
        
        if file_type not in self.supported_formats:
            raise Exception(f"不支持的文件类型: {file_type}")
        
        if not self.supported_formats[file_type]:
            raise Exception(f"处理{file_type}文件的库未安装")
        
        if file_type in ['pptx', 'ppt']:
            return self.extract_tables_from_pptx(file_path)
        elif file_type in ['docx', 'doc']:
            return self.extract_tables_from_docx(file_path)
        elif file_type == 'pdf':
            return self.extract_tables_from_pdf(file_path)
        else:
            raise Exception(f"不支持的文件类型: {file_type}")
    
    def save_tables_to_excel(self, tables: List[Dict[str, Any]], output_path: str) -> Dict[str, Any]:
        """将提取的表格保存到Excel文件"""
        try:
            wb = Workbook()
            
            # 删除默认工作表
            wb.remove(wb.active)
            
            results = {
                'total_tables': len(tables),
                'sheets_created': [],
                'file_path': output_path
            }
            
            for i, table_info in enumerate(tables, 1):
                # 创建工作表
                sheet_name = f"Table_{i}"
                if len(sheet_name) > 31:  # Excel工作表名称限制
                    sheet_name = f"T{i}"
                
                ws = wb.create_sheet(title=sheet_name)
                
                # 添加表格数据
                table_data = table_info['data']
                for row_idx, row_data in enumerate(table_data, 1):
                    for col_idx, cell_value in enumerate(row_data, 1):
                        ws.cell(row=row_idx, column=col_idx, value=cell_value)
                
                # 记录工作表信息
                sheet_info = {
                    'sheet_name': sheet_name,
                    'table_info': table_info,
                    'rows': len(table_data),
                    'columns': len(table_data[0]) if table_data else 0
                }
                results['sheets_created'].append(sheet_info)
            
            # 保存文件
            wb.save(output_path)
            
            return results
            
        except Exception as e:
            logger.error(f"保存Excel文件失败: {e}")
            raise Exception(f"保存Excel文件失败: {str(e)}")
    
    def extract_and_save(self, url: str, output_filename: str = None) -> Dict[str, Any]:
        """从URL下载文档，提取表格并保存为Excel"""
        temp_file = None
        try:
            # 下载文件
            temp_file = self.download_file(url)
            
            # 确定文件类型
            file_type = os.path.splitext(temp_file)[1].lower().lstrip('.')
            
            # 提取表格
            tables = self.extract_tables_from_document(temp_file, file_type)
            
            if not tables:
                return {
                    'success': False,
                    'message': '未在文档中找到表格',
                    'total_tables': 0
                }
            
            # 生成输出文件名
            if not output_filename:
                base_name = os.path.splitext(os.path.basename(temp_file))[0]
                output_filename = f"{base_name}_extracted_tables.xlsx"
            
            output_path = os.path.join(tempfile.gettempdir(), output_filename)
            
            # 保存到Excel
            save_result = self.save_tables_to_excel(tables, output_path)
            
            return {
                'success': True,
                'message': f'成功提取{len(tables)}个表格',
                'total_tables': len(tables),
                'output_file': output_path,
                'tables_info': tables,
                'save_result': save_result
            }
            
        except Exception as e:
            logger.error(f"提取表格失败: {e}")
            return {
                'success': False,
                'message': f'提取表格失败: {str(e)}',
                'error': str(e)
            }
        finally:
            # 清理临时文件
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except:
                    pass

def extract_tables_from_document_url(url: str, output_filename: str = None) -> Dict[str, Any]:
    """从文档URL提取表格的便捷函数"""
    extractor = DocumentExtractor()
    return extractor.extract_and_save(url, output_filename)
